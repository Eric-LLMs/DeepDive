"""Unit tests for the document-generation tools (doc_outline / doc_mindmap / doc_slides).

The tools read a drive asset's bytes and save generated .pptx files back into the drive.
We exercise them with hand-rolled fakes (same style as ``_drive_fakes.py``): the storage
is a real :class:`LocalStorage` at tmp_path and the SQL repositories are monkeypatched so
no database is touched.
"""
from __future__ import annotations

from pathlib import Path
from uuid import uuid4

import pytest

from agent.engine.decisions import ToolExecution
from agent.engine.runtime import ToolRuntime
from agent import Context
from core.infrastructure.request_context import set_request_user
from core.infrastructure.storage import LocalStorage, object_key

from apps.api.tools import document_tools
from tests._drive_fakes import FakeAssets, FakeObjects


class _FakeLLM:
    def __init__(self) -> None:
        self.calls: list[tuple[str, str]] = []

    async def complete(self, text: str, system: str) -> str:
        self.calls.append((text, system))
        if "structure analyst" in system:
            return "# Overview\n## Details"
        if "slide deck writer" in system:
            return (
                '[{"title": "Intro", "bullets": ["Point A", "Point B"]},'
                '{"title": "Conclusion", "bullets": ["Takeaway"]}]'
            )
        return "central topic\n\tbranch one\n\t\tleaf"

    async def chat(self, *_args, **_kwargs):  # pragma: no cover - unused surface
        return None


def _ctx(tmp_path: Path) -> Context:
    ctx = Context()
    ctx.provide("storage", LocalStorage(tmp_path))
    ctx.provide("session_factory", object)  # unused: repos are monkeypatched
    return ctx


@pytest.mark.parametrize(
    ("raw", "expected"),
    [
        ('[{"title": "A", "bullets": ["x", "y"]}]', [("A", "x\ny")]),
        ('```json\n[{"title": "A", "bullets": ["x"]}]\n```', [("A", "x")]),
        ("prefix [{\"title\": \"A\", \"bullets\": []}] suffix", [("A", "")]),
        ("not json at all", []),
        ("[{\"bullets\": [\"x\"]}]", []),  # missing title
    ],
)
def test_parse_slides_json(raw, expected):
    assert document_tools._parse_slides_json(raw) == expected


def test_build_text_pptx(tmp_path):
    from pptx import Presentation

    out = tmp_path / "deck.pptx"
    document_tools.media_lib.build_text_pptx(
        [("Intro", "Point A\nPoint B"), ("Conclusion", "Takeaway")],
        out,
        title="My Deck",
    )
    prs = Presentation(str(out))
    # title slide + 2 content slides
    assert len(prs.slides._sldIdLst) == 3
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    assert any("My Deck" in t for t in texts)
    assert any("Intro" in t for t in texts)


async def test_doc_outline_executes(monkeypatch, tmp_path):
    set_request_user(uuid4())
    assets = FakeAssets()
    asset = await assets.create(uuid4(), "notes.md", object_sha256="a" * 64)
    monkeypatch.setattr(document_tools, "SqlAssetRepository", lambda sf: assets)
    monkeypatch.setattr(
        document_tools, "_load_asset_bytes",
        lambda asset_id, ctx: b"# Title\n\nSome body text.",
    )

    runtime = ToolRuntime()
    document_tools.register(runtime, _ctx(tmp_path), _FakeLLM())

    res = await runtime.execute(ToolExecution("c1", "doc_outline", {"asset_id": str(asset.id)}))
    assert res.is_error is False
    assert "# Overview" in str(res.value)


async def test_doc_slides_saves_pptx_to_drive(monkeypatch, tmp_path):
    user_id = uuid4()
    set_request_user(user_id)
    assets = FakeAssets()
    asset = await assets.create(user_id, "notes.md", object_sha256="b" * 64)
    objects = FakeObjects()
    monkeypatch.setattr(document_tools, "SqlAssetRepository", lambda sf: assets)
    monkeypatch.setattr(document_tools, "SqlGlobalObjectRepository", lambda sf: objects)
    monkeypatch.setattr(
        document_tools, "_load_asset_bytes",
        lambda asset_id, ctx: b"# Title\n\nSome body text.",
    )

    ctx = _ctx(tmp_path)
    runtime = ToolRuntime()
    document_tools.register(runtime, ctx, _FakeLLM())

    res = await runtime.execute(ToolExecution("c2", "doc_slides", {"asset_id": str(asset.id), "count": 5}))
    assert res.is_error is False
    assert "Saved" in str(res.value)

    # A new READY asset was created for the generated deck.
    saved = [a for a in assets.rows.values() if a.name.startswith("notes_slides_") and a.name.endswith(".pptx")]
    assert len(saved) == 1
    deck = saved[0]
    assert deck.file_status == "READY"
    assert deck.rag_status == "NOT_STARTED"
    assert deck.object_sha256 is not None

    # Its bytes are physically stored, and the object row is ref-counted.
    stored = await ctx.resolve("storage").get(object_key(deck.object_sha256))
    assert stored and stored.startswith(b"PK")  # a valid zip/pptx
    assert objects.rows[deck.object_sha256].ref_count >= 1


async def test_doc_slides_no_extractable_text(monkeypatch, tmp_path):
    set_request_user(uuid4())
    assets = FakeAssets()
    asset = await assets.create(uuid4(), "video.mp4", object_sha256="c" * 64)
    monkeypatch.setattr(document_tools, "SqlAssetRepository", lambda sf: assets)
    monkeypatch.setattr(
        document_tools, "_load_asset_bytes", lambda asset_id, ctx: b"\x00\x01",  # not text
    )

    runtime = ToolRuntime()
    document_tools.register(runtime, _ctx(tmp_path), _FakeLLM())
    res = await runtime.execute(ToolExecution("c3", "doc_slides", {"asset_id": str(asset.id)}))
    assert res.is_error is False
    assert "No extractable text" in str(res.value)
