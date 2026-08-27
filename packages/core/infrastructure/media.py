"""Media utilities for the desktop workbench ingestion pipeline.

Pure-stdlib subtitle parsing (SRT / VTT / LRC) plus optional, lazily-imported
helpers for keyframe extraction (ffmpeg) and PPT / PDF "book" generation. The
heavy dependencies (python-pptx, reportlab, imageio-ffmpeg) are imported inside
the functions that need them, so this module stays importable without them.
"""
from __future__ import annotations

import re
import subprocess
from dataclasses import dataclass
from pathlib import Path


@dataclass
class SubtitleCue:
    index: int
    start_ms: int
    end_ms: int
    text: str


# ── time helpers ──
def _hms_to_ms(h: int, m: int, s: int, ms: int) -> int:
    return h * 3600000 + m * 60000 + s * 1000 + ms


def _srt_ts_to_ms(ts: str) -> int:
    h, m, rest = ts.split(":")
    s, ms = rest.split(",")
    return _hms_to_ms(int(h), int(m), int(s), int(ms.ljust(3, "0")[:3]))


def _read_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8-sig")
    except UnicodeDecodeError:
        return path.read_text(encoding="latin-1")


# ── SRT ──
_SRT_RE = re.compile(
    r"(\d+)\s*?\n(\d{1,2}:\d{2}:\d{2},\d{1,3})\s*-->\s*(\d{1,2}:\d{2}:\d{2},\d{1,3})\s*?\n(.*?)(?=\n\s*\n|\Z)",
    re.DOTALL,
)


def parse_srt_text(text: str) -> list[SubtitleCue]:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    cues: list[SubtitleCue] = []
    for m in _SRT_RE.finditer(text):
        idx, start, end, body = m.groups()
        clean = " ".join(body.split())
        if clean:
            cues.append(SubtitleCue(int(idx), _srt_ts_to_ms(start), _srt_ts_to_ms(end), clean))
    return cues


# ── VTT ──
_VTT_TS_RE = re.compile(
    r"(?:(\d+):)?(\d{2}):(\d{2})\.(\d{1,3})\s*-->\s*(?:(\d+):)?(\d{2}):(\d{2})\.(\d{1,3})"
)


def _strip_vtt_tags(line: str) -> str:
    line = re.sub(r"<\d+:\d{2}:\d{2}\.\d{1,3}>", "", line)
    line = re.sub(r"</?c(?:\.\w+)*>", "", line)
    line = re.sub(r"</?v(?:\.\w+)*>", "", line)
    line = re.sub(r"</?[a-z]+>", "", line)
    line = line.replace("\\N", " ").replace("\\n", " ").replace("\\h", "")
    return re.sub(r" {2,}", " ", line).strip()


def _vtt_groups_to_ms(m: re.Match, base: int) -> int:
    h = int(m.group(base) or 0)
    mi = int(m.group(base + 1))
    s = int(m.group(base + 2))
    ms = int(m.group(base + 3).ljust(3, "0")[:3])
    return _hms_to_ms(h, mi, s, ms)


def parse_vtt_text(text: str) -> list[SubtitleCue]:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    body = text
    if body.lstrip().startswith("WEBVTT"):
        idx = body.find("\n\n")
        if idx != -1:
            body = body[idx + 2 :]
        else:
            body = body[body.find("\n") + 1 :]

    cues: list[SubtitleCue] = []
    for block in re.split(r"\n{2,}", body):
        ts_line = ""
        text_lines: list[str] = []
        for line in block.splitlines():
            s = line.strip()
            if not s:
                continue
            if "-->" in s:
                ts_line = s
            elif not s.startswith(("NOTE", "STYLE", "REGION")):
                clean = _strip_vtt_tags(s)
                if clean:
                    text_lines.append(clean)
        if not ts_line or not text_lines:
            continue
        m = _VTT_TS_RE.search(ts_line)
        if not m:
            continue
        cues.append(
            SubtitleCue(
                len(cues) + 1,
                _vtt_groups_to_ms(m, 1),
                _vtt_groups_to_ms(m, 5),
                " ".join(text_lines),
            )
        )
    return cues


# ── LRC ──
_LRC_RE = re.compile(r"\[(\d+):(\d+)\.(\d+)\](.*)")


def parse_lrc_text(text: str) -> list[SubtitleCue]:
    cues: list[SubtitleCue] = []
    for line in text.replace("\r\n", "\n").replace("\r", "\n").splitlines():
        m = _LRC_RE.match(line.strip())
        if not m:
            continue
        mi, s, ms, txt = m.groups()
        cues.append(
            SubtitleCue(
                len(cues) + 1,
                _hms_to_ms(0, int(mi), int(s), int(ms.ljust(3, "0")[:3])),
                0,
                txt.strip(),
            )
        )
    # Back-fill end times from the next cue's start.
    for i, cue in enumerate(cues):
        cue.end_ms = cues[i + 1].start_ms if i + 1 < len(cues) else cue.start_ms + 3000
    return [c for c in cues if c.text]


def parse_subtitles(path: str | Path) -> list[SubtitleCue]:
    """Parse a subtitle file, dispatching on extension (.srt/.vtt/.lrc)."""
    path = Path(path)
    text = _read_text(path)
    ext = path.suffix.lower()
    if ext == ".vtt":
        return parse_vtt_text(text)
    if ext == ".lrc":
        return parse_lrc_text(text)
    return parse_srt_text(text)


# ── ffmpeg keyframe extraction ──
def _ffmpeg_exe() -> str:
    import imageio_ffmpeg

    return imageio_ffmpeg.get_ffmpeg_exe()


def video_duration_ms(video_path: str | Path) -> int:
    """Return video duration in milliseconds (0 if it can't be determined)."""
    proc = subprocess.run(
        [_ffmpeg_exe(), "-i", str(video_path)],
        capture_output=True,
        text=True,
    )
    m = re.search(r"Duration:\s*(\d+):(\d+):(\d+\.\d+)", proc.stderr or "")
    if not m:
        return 0
    return _hms_to_ms(int(m.group(1)), int(m.group(2)), int(m.group(3)), 0)


def extract_keyframes(
    video_path: str | Path,
    timestamps_ms: list[int],
    out_dir: str | Path,
    interval_ms: int = 5000,
) -> list[str]:
    """Extract one frame per timestamp; if none given, sample every ``interval_ms``.

    Returns the list of generated PNG paths (same order as the timestamps).
    """
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    ffmpeg = _ffmpeg_exe()

    if not timestamps_ms:
        duration = video_duration_ms(video_path)
        timestamps_ms = list(range(0, duration, interval_ms)) or [0]

    paths: list[str] = []
    for i, ms in enumerate(timestamps_ms):
        out = out_dir / f"frame_{i:04d}.png"
        subprocess.run(
            [
                ffmpeg,
                "-y",
                "-ss",
                f"{ms / 1000.0:.3f}",
                "-i",
                str(video_path),
                "-frames:v",
                "1",
                str(out),
            ],
            capture_output=True,
            check=True,
        )
        if out.exists():
            paths.append(str(out))
    return paths


# ── PPT generation ──
def build_pptx(
    slides: list[tuple[str, str]],
    out_path: str | Path,
    title: str = "",
) -> str:
    """Build a .pptx where each slide is (image_path, text)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    if title:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.3), Inches(1.5))
        para = box.text_frame.paragraphs[0]
        para.text = title
        para.font.size = Pt(40)
        para.font.bold = True

    for image, text in slides:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(image, Inches(0.5), Inches(0.5), height=Inches(5.5))
        box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.3), Inches(1.1))
        box.text_frame.word_wrap = True
        box.text_frame.text = text

    out_path = Path(out_path)
    prs.save(str(out_path))
    return str(out_path)


def build_text_pptx(
    slides: list[tuple[str, str]],
    out_path: str | Path,
    title: str = "",
) -> str:
    """Build a text-only .pptx: one bullet slide per ``(heading, bullets)`` tuple.

    Unlike :func:`build_pptx` (which stamps a frame image per page), this variant is
    for document-derived decks — the agent turns extracted document text into a slide
    outline and we lay it out as title + bullet text boxes. ``bullets`` may contain
    newline-separated points; each becomes its own paragraph.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    if title:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.3), Inches(1.5))
        para = box.text_frame.paragraphs[0]
        para.text = title
        para.font.size = Pt(40)
        para.font.bold = True

    for heading, bullets in slides:
        slide = prs.slides.add_slide(blank)
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11.9), Inches(1.1))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = heading
        tp.font.size = Pt(32)
        tp.font.bold = True

        body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(11.9), Inches(5.2))
        tf = body_box.text_frame
        tf.word_wrap = True
        lines = [ln for ln in (bullets or "").split("\n") if ln.strip()]
        for i, line in enumerate(lines):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = f"• {line.strip()}"
            para.font.size = Pt(18)
            para.space_after = Pt(8)

    out_path = Path(out_path)
    prs.save(str(out_path))
    return str(out_path)


# ── PDF "book" generation ──
def _wrap_text(text: str, width_fn, max_width: int) -> list[str]:
    lines: list[str] = []
    current = ""
    for ch in text:
        if width_fn(current + ch) > max_width:
            if current:
                lines.append(current)
            current = ch
        else:
            current += ch
    if current:
        lines.append(current)
    return lines


def build_pdf(
    slides: list[tuple[str, str]],
    out_path: str | Path,
    title: str = "",
) -> str:
    """Build a PDF "book" where each page is (image, text)."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.pdfgen import canvas

    # STSong-Light is a built-in CID font that renders CJK text.
    pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))

    page_w, page_h = A4
    out_path = Path(out_path)
    c = canvas.Canvas(str(out_path), pagesize=A4)
    font = "STSong-Light"
    font_size = 12
    c.setFont(font, font_size)
    margin = 0.75 * inch
    text_w = page_w - 2 * margin
    img_h = (page_h - 2 * margin) * 0.55
    text_y = margin + 0.4 * inch

    if title:
        c.setFont(font, 24)
        c.drawCentredString(page_w / 2, page_h / 2, title)
        c.showPage()
        c.setFont(font, font_size)

    for image, text in slides:
        try:
            c.drawImage(image, margin, page_h - margin - img_h, width=text_w, height=img_h,
                        preserveAspectRatio=True, anchor="n")
        except Exception:  # noqa: BLE001 - a bad frame should not abort the book
            pass
        c.setFont(font, font_size)
        for i, line in enumerate(_wrap_text(text, lambda s: c.stringWidth(s, font, font_size), int(text_w))):
            c.drawString(margin, text_y - i * 16, line)
        c.showPage()

    c.save()
    return str(out_path)
